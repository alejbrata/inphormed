#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os, csv, random, argparse, textwrap, datetime, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt

random.seed(42)

# ---------- Referencias reales (trazables) ----------
REFS = [
    {
        "key": "PIONEER_NEJM_2016",
        "type": "doi",
        "id": "10.1056/NEJMoa1504370",
        "url": "https://www.nejm.org/doi/full/10.1056/NEJMoa1504370",
        "facts": [
            "Adalimumab mejora la respuesta clínica (HiSCR) a semana 12 frente a placebo en HS.",
            "La variable primaria incluyó HiSCR en PIONEER I/II.",
        ],
    },
    {
        "key": "HiSCR_JEADV_2016",
        "type": "pmid",
        "id": "26201313",
        "url": "https://pubmed.ncbi.nlm.nih.gov/26201313/",
        "facts": [
            "HiSCR: ≥50% reducción en abscesos y nódulos inflamatorios, sin incremento en abscesos o fístulas drenantes.",
        ],
    },
    {
        "key": "Hurley_StatPearls",
        "type": "url",
        "id": "NBK534867",
        "url": "https://www.ncbi.nlm.nih.gov/books/NBK534867/",
        "facts": [
            "Hurley I: abscesos sin túneles ni cicatrices; Hurley III: tractos fistulosos interconectados y afectación extensa.",
        ],
    },
    {
        "key": "HS_Review_PMC6330680",
        "type": "url",
        "id": "PMC6330680",
        "url": "https://pmc.ncbi.nlm.nih.gov/articles/PMC6330680/",
        "facts": [
            "La HS es crónica, inflamatoria y recidivante; el Hurley clasifica severidad en I–III.",
        ],
    },
    {
        "key": "Secukinumab_FDA_2023",
        "type": "url",
        "id": "COSENTYX_FDA_2023",
        "url": "https://www.novartis.com/news/media-releases/fda-approves-novartis-cosentyx-first-new-biologic-treatment-option-hidradenitis-suppurativa-patients-nearly-decade",
        "facts": [
            "En 2023, la FDA aprobó secukinumab para HS moderada-grave en adultos.",
        ],
    },
    {
        "key": "Bimekizumab_FDA_2024",
        "type": "url",
        "id": "BIMZELX_FDA_2024",
        "url": "https://www.accessdata.fda.gov/drugsatfda_docs/label/2024/761151s010lbl.pdf",
        "facts": [
            "Bimekizumab (Bimzelx) tiene indicación aprobada en HS moderada-grave en adultos (EE. UU., 2024).",
        ],
    },
    {
        "key": "Bimekizumab_EU_2024",
        "type": "url",
        "id": "BIMZELX_EU_2024",
        "url": "https://www.ema.europa.eu/en/medicines/human/EPAR/bimzelx",
        "facts": [
            "La UE aprobó bimekizumab para HS moderada-grave (EC 2024; ver EPAR).",
        ],
    },
    {
        "key": "S2k_JDV_2025",
        "type": "doi",
        "id": "10.1111/jdv.20472",
        "url": "https://onlinelibrary.wiley.com/doi/10.1111/jdv.20472",
        "facts": [
            "Guías europeas S2k (2025) con recomendaciones de manejo por gravedad y líneas de tratamiento.",
        ],
    },
]

# ---------- Plantillas de claims falsos (rojo) ----------
FALSE_CLAIMS = [
    "La hidradenitis supurativa se debe principalmente a una mala higiene personal.",  # falso
    "HiSCR requiere una reducción del 100% de las lesiones para considerarse respondedor.",  # falso
    "Bimekizumab actúa inhibiendo IL-23 como mecanismo principal en HS.",  # falso (es IL-17A/F)
    "Adalimumab cura definitivamente la HS en el 90% de los pacientes en 12 semanas.",  # falso
    "Secukinumab reduce el riesgo de melanoma en HS en un 50% a las 8 semanas.",  # inventado
]

DISCLAIMER = (
    "DOCUMENTO SINTÉTICO PARA PRUEBAS DEL TFM — NO USO CLÍNICO. "
    "Contiene afirmaciones VERDES (ciertas/similares a la fuente), AMARILLAS (dudosas/ambiguas) y ROJAS (falsas). "
    "Propósito: evaluar trazabilidad y detección automática. "
)

def pick_green():
    ref = random.choice(REFS)
    claim = random.choice(ref["facts"])
    return claim, ref

def pick_yellow():
    # Toma un hecho real y lo vuelve ambiguo / parcial (sin dejar de ser plausible)
    ref = random.choice(REFS)
    base = random.choice(ref["facts"])
    # heurística: suavizar/ambigüedad
    tweak = re.sub(r"\b(\d+%|\d+\s*semanas?)\b", "algunas semanas", base)
    tweak = tweak.replace("aprob", "autoriz")  # menos preciso
    tweak = tweak.replace("≥50%", "alrededor del 50%")
    return tweak, ref

def pick_red():
    claim = random.choice(FALSE_CLAIMS)
    # Usa una referencia real desparejada (para que el validador lo marque rojo)
    ref = random.choice(REFS)
    return claim, ref

def mk_dirs(root):
    Path(root, "pptx").mkdir(parents=True, exist_ok=True)
    Path(root, "docx").mkdir(parents=True, exist_ok=True)

def add_ppt_slide(prs, title, body):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.LEFT

def color_rgb(label):
    return {"GREEN": RGBColor(0,150,0), "YELLOW": RGBColor(255,165,0), "RED": RGBColor(200,0,0)}[label]

def build_pptx(path, items):
    prs = Presentation()
    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hidradenitis supurativa — Claims sintéticos"
    slide.placeholders[1].text = DISCLAIMER

    # Slides por claim (2 claims por slide)
    for i, it in enumerate(items, start=1):
        title = f"Claim {i} — {it['label']}"
        body = f"{it['text']}\n\nCita: {it['cite']} ({it['url']})"
        add_ppt_slide(prs, title, body)

        # Colorear el título según etiqueta
        for shape in prs.slides[-1].shapes:
            if hasattr(shape, "text_frame") and shape.text == title:
                for r in shape.text_frame.paragraphs[0].runs:
                    r.font.color.rgb = color_rgb(it['label'])

    prs.save(path)

def build_docx(path, items):
    doc = Document()
    p = doc.add_paragraph(DISCLAIMER)
    p.runs[0].bold = True
    p.runs[0].font.size = DocxPt(10)

    doc.add_heading("Hidradenitis supurativa — Claims sintéticos", level=1)

    for i, it in enumerate(items, start=1):
        h = doc.add_heading(f"Claim {i} — {it['label']}", level=2)
        # coloreo simple por etiqueta (no todos los visores respetan highlight)
        color_map = {"GREEN": 3, "YELLOW": 7, "RED": 6}  # bright-green, yellow, magenta
        for run in h.runs:
            run.font.highlight_color = color_map[it['label']]
        para = doc.add_paragraph(it["text"])
        doc.add_paragraph(f"Cita: {it['cite']} ({it['url']})").italic = True

    doc.save(path)

def make_items(n_claims=6, p_green=0.4, p_yellow=0.3):
    items = []
    for _ in range(n_claims):
        r = random.random()
        if r < p_green:
            txt, ref = pick_green(); label = "GREEN"
        elif r < p_green + p_yellow:
            txt, ref = pick_yellow(); label = "YELLOW"
        else:
            txt, ref = pick_red(); label = "RED"

        cite = f"{ref['type'].upper()}:{ref['id']}" if ref["type"] in ("doi","pmid") else ref["url"]
        items.append({
            "label": label,
            "text": txt,
            "cite": cite,
            "url": ref["url"],
            "ref_key": ref["key"],
        })
    return items

def main(out_root, n_pptx=50, n_docx=50, claims_per_doc=6):
    ts = datetime.datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    out_root = Path(out_root)
    mk_dirs(out_root)
    gt_path = out_root / "ground_truth.csv"

    with open(gt_path, "w", newline="", encoding="utf-8") as f:
        wr = csv.writer(f)
        wr.writerow(["file","type","idx","label","text","citation","url","ref_key"])

        # PPTX
        for i in range(1, n_pptx+1):
            items = make_items(n_claims=claims_per_doc)
            fname = f"hs_synth_{ts}_{i:03d}.pptx"
            fpath = out_root / "pptx" / fname
            build_pptx(str(fpath), items)
            for j, it in enumerate(items, start=1):
                wr.writerow([fname,"pptx",j,it["label"],it["text"],it["cite"],it["url"],it["ref_key"]])

        # DOCX
        for i in range(1, n_docx+1):
            items = make_items(n_claims=claims_per_doc)
            fname = f"hs_synth_{ts}_{i:03d}.docx"
            fpath = out_root / "docx" / fname
            build_docx(str(fpath), items)
            for j, it in enumerate(items, start=1):
                wr.writerow([fname,"docx",j,it["label"],it["text"],it["cite"],it["url"],it["ref_key"]])

    print(f"Generado dataset en: {out_root}\n- PPTX: {n_pptx}\n- DOCX: {n_docx}\n- GT: {gt_path}")

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=str, default="synthetic/hs")
    ap.add_argument("--n_pptx", type=int, default=50)
    ap.add_argument("--n_docx", type=int, default=50)
    ap.add_argument("--claims_per_doc", type=int, default=6)
    args = ap.parse_args()
    main(args.out, args.n_pptx, args.n_docx, args.claims_per_doc)
