import os
import re
import uuid
from typing import List, Dict, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

from docx import Document
from docx.enum.text import WD_COLOR_INDEX

COLOR_RGB = {
    "green":  RGBColor(0, 150, 0),
    "yellow": RGBColor(255, 165, 0),
    "red":    RGBColor(200, 0, 0),
}

DOCX_HILITE = {
    "green":  WD_COLOR_INDEX.BRIGHT_GREEN,
    "yellow": WD_COLOR_INDEX.YELLOW,
    "red":    WD_COLOR_INDEX.PINK,  # no hay ‘orange’, usamos rosa como alerta
}

def _build_colors(color: str) -> Tuple[RGBColor, WD_COLOR_INDEX]:
    return COLOR_RGB.get(color, COLOR_RGB["red"]), DOCX_HILITE.get(color, WD_COLOR_INDEX.RED)

def write_snippets_html(base_uid: str, findings: List[Dict]) -> str:
    os.makedirs("outputs/snippets", exist_ok=True)
    path = f"outputs/snippets/snippets_{base_uid}.html"
    rows = []
    for idx, f in enumerate(findings, start=1):
        rows.append(f"""
        <section id="claim-{idx}" style="margin:1rem 0;padding:1rem;border:1px solid #ddd;border-radius:10px">
          <h3>Claim {idx} — {f['color'].upper()} (score={int(f['score'])})</h3>
          <p><b>Texto del slide/bloque:</b><br>{escape_html(f['claim'])}</p>
          <p><b>Extracto fuente:</b><br><pre style="white-space:pre-wrap">{escape_html(f.get('source_excerpt') or '')}</pre></p>
          <p><a href="{f.get('source_url','')}" target="_blank">Abrir paper</a></p>
        </section>
        """)
    html = f"""<!doctype html><html lang="es"><meta charset="utf-8">
      <title>Snippets Inphormed</title>
      <body style="font-family:Arial,sans-serif;max-width:900px;margin:2rem auto">
      <h1>Snippets de evidencias — Inphormed</h1>
      {''.join(rows)}
      </body></html>"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(html)
    return path

def escape_html(s: str) -> str:
    return (s or "").replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

# ---------------- PPTX ----------------

def annotate_pptx(original_path: str, out_path: str, findings: List[Dict]) -> str:
    prs = Presentation(original_path)
    rgb_green, _ = _build_colors("green")  # default
    # mapeo slide_index -> lista de claims para ese slide (orden)
    claims_by_slide = {}
    for f in findings:
        claims_by_slide.setdefault(f["slide"], []).append(f)

    for i, slide in enumerate(prs.slides, start=1):
        slide_claims = claims_by_slide.get(i, [])
        if not slide_claims:
            continue

        # Por cada shape con text_frame, coloreamos si el texto contiene parte del claim
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    text_lower = run.text.strip().lower()
                    if not text_lower:
                        continue
                    # buscamos el claim más similar asignado a este slide (heurística simple)
                    for idx, f in enumerate(slide_claims, start=1):
                        color = f["color"]
                        rgb, _ = _build_colors(color)
                        # heurística: si comparte n-gramas clave, colorea y enlaza
                        if _weak_match(text_lower, f["claim"].lower()):
                            run.font.color.rgb = rgb
                            run.font.bold = True
                            run.font.size = Pt(para.font.size.pt if para.font.size else 14)
                            # añade hyperlink al snippet local (preferente) o al paper
                            url = f.get("snippet_url") or f.get("source_url")
                            if url:
                                try:
                                    hlink = run.hyperlink
                                    hlink.address = url
                                except Exception:
                                    pass
    prs.save(out_path)
    return out_path

def _weak_match(text: str, claim: str) -> bool:
    # match ligero: comparte ≥3 palabras de 6+ letras
    words_t = {w for w in re.findall(r"[a-záéíóúñ]{6,}", text) }
    words_c = {w for w in re.findall(r"[a-záéíóúñ]{6,}", claim)}
    return len(words_t.intersection(words_c)) >= 3

# ---------------- DOCX ----------------

def annotate_docx(original_path: str, out_path: str, findings: List[Dict]) -> str:
    doc = Document(original_path)
    for para in doc.paragraphs:
        ptext = para.text.strip().lower()
        if not ptext:
            continue
        for idx, f in enumerate(findings, start=1):
            claim = (f["claim"] or "").lower()
            if not _weak_match(ptext, claim):
                continue
            _, hi = _build_colors(f["color"])
            # coloreo simple a nivel de párrafo (más robusto que partir runs arbitrariamente)
            for run in para.runs:
                run.font.highlight_color = hi
                run.font.bold = True
            # añade hyperlink al final del párrafo (fallback)
            url = f.get("snippet_url") or f.get("source_url")
            if url:
                para.add_run(" [ver evidencia]").font.underline = True
                try:
                    # truco: docx no tiene API directa para hyperlink en run; aceptamos subrayado + URL visible
                    para.add_run(f" ({url})")
                except Exception:
                    pass
    doc.save(out_path)
    return out_path
