# streamlit_app.py
import os
import json
from typing import List, Dict, Any
import streamlit as st

# ─────────────────────────────────────────────────────────────
# Config de conexión (ajusta por ENV si no es localhost:8000)
# ─────────────────────────────────────────────────────────────
BACKEND_URL = os.getenv("BACKEND_URL", "http://localhost:8000")
CHAT_ENDPOINT = os.getenv("CHAT_ENDPOINT", "/api/chat")
CHAT_URL = BACKEND_URL.rstrip("/") + CHAT_ENDPOINT

st.set_page_config(page_title="Inphormed", layout="wide")

# ─────────────────────────────────────────────────────────────
# Helper HTTP (usa requests si está; si no, urllib)
# ─────────────────────────────────────────────────────────────
def post_json(url: str, payload: Dict[str, Any]) -> Dict[str, Any]:
    try:
        import requests
        r = requests.post(url, json=payload, timeout=60)
        r.raise_for_status()
        return r.json()
    except Exception:
        import urllib.request, urllib.error
        data = json.dumps(payload).encode("utf-8")
        req = urllib.request.Request(url, data=data, headers={"Content-Type": "application/json"})
        with urllib.request.urlopen(req, timeout=60) as resp:
            return json.loads(resp.read().decode("utf-8"))

# ─────────────────────────────────────────────────────────────
# Utilidades mínimas para el validador (DOCX/PPTX/TXT)
# ─────────────────────────────────────────────────────────────
def parse_docx(file_bytes: bytes) -> List[str]:
    try:
        from docx import Document
        import io
        doc = Document(io.BytesIO(file_bytes))
        lines = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                lines.append(t)
        return lines
    except Exception as e:
        return [f"[ERROR DOCX] {e}"]

def parse_pptx(file_bytes: bytes) -> List[str]:
    try:
        from pptx import Presentation
        import io
        pres = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    txt = (shape.text_frame.text or "").strip()
                    if txt:
                        lines.append(txt)
        return lines
    except Exception as e:
        return [f"[ERROR PPTX] {e}"]

def dedupe_keep_order(items: List[str]) -> List[str]:
    seen, out = set(), []
    for x in items:
        if x and x not in seen:
            seen.add(x); out.append(x)
    return out

# ─────────────────────────────────────────────────────────────
# VISTAS
# ─────────────────────────────────────────────────────────────
def vista_chatbot_hs():
    st.header("Chatbot – Hidradenitis supurativa")

    # historial en sesión (solo 'user' y 'assistant'; el 'system' lo pone el backend)
    if "chat_hs" not in st.session_state:
        st.session_state.chat_hs = [
            {"role": "assistant", "content": "Hola, ¿en qué te ayudo?"}
        ]

    # Render del histórico
    for m in st.session_state.chat_hs:
        st.chat_message(m["role"]).markdown(m["content"])

    # Input del usuario
    user_msg = st.chat_input("Pregunta o describe tu caso…")
    if user_msg:
        st.session_state.chat_hs.append({"role": "user", "content": user_msg})
        st.chat_message("user").markdown(user_msg)

        # Construye el payload para TU backend (SIN 'system')
        messages_for_api: List[Dict[str, str]] = [
            {"role": m["role"], "content": m["content"]}
            for m in st.session_state.chat_hs
            if m["role"] in ("user", "assistant")
        ]
        payload = {"messages": messages_for_api, "topic": "hidradenitis supurativa"}

        try:
            with st.spinner("Pensando…"):
                # ***** AQUÍ se hace el POST al backend *****
                resp = post_json(CHAT_URL, payload)

            # Tu ruta devuelve {"reply": "..."}
            reply = (resp or {}).get("reply", "").strip() or "No recibí contenido del modelo."
            st.session_state.chat_hs.append({"role": "assistant", "content": reply})
            st.chat_message("assistant").markdown(reply)
        except Exception as e:
            err = f"Error llamando a {CHAT_URL}: {e}"
            st.error(err)
            st.session_state.chat_hs.append({"role": "assistant", "content": err})

def vista_validador_claims():
    st.header("Validador de claims (UI mínima)")
    st.caption("De momento: lectura de claims. Siguiente paso: llamada al validador con IA/RAG.")

    tabs = st.tabs(["Pegar texto", "Subir archivo (DOCX / PPTX / TXT)"])

    with tabs[0]:
        st.subheader("Pegar claims (uno por línea)")
        txt = st.text_area(
            "Claims",
            height=200,
            placeholder="Mejora la supervivencia global…\nReducción del 30% frente a SOC…",
        )
        if st.button("Previsualizar claims (texto)"):
            claims = dedupe_keep_order([line.strip() for line in txt.splitlines() if line.strip()])
            if not claims:
                st.warning("Escribe al menos un claim.")
            else:
                st.success(f"{len(claims)} claims detectados (texto).")
                st.dataframe({"claim": claims}, use_container_width=True)

    with tabs[1]:
        st.subheader("Subir archivo")
        f = st.file_uploader("DOCX / PPTX / TXT", type=["docx","pptx","txt"])
        if st.button("Previsualizar claims (archivo)", disabled=not f):
            if not f:
                st.warning("Selecciona un archivo.")
            else:
                ext = (f.name.split(".")[-1] or "").lower()
                data = f.read()
                if ext == "docx":
                    claims = parse_docx(data)
                elif ext == "pptx":
                    claims = parse_pptx(data)
                elif ext == "txt":
                    claims = [line.strip() for line in data.decode("utf-8", errors="ignore").splitlines() if line.strip()]
                else:
                    claims = []

                claims = [c for c in claims if c and not c.startswith("[ERROR")]
                claims = dedupe_keep_order(claims)

                if not claims:
                    st.error("No se encontraron claims en el archivo.")
                else:
                    st.success(f"{len(claims)} claims detectados (archivo).")
                    st.dataframe({"claim": claims}, use_container_width=True)

def vista_generador_material():
    st.header("Generador de material")
    st.info("Placeholder limpio. Aquí añadiremos plantillas y salida (PDF/PowerPoint) cuando toque.")

# ─────────────────────────────────────────────────────────────
# Router (tus 3 entradas)
# ─────────────────────────────────────────────────────────────
PAGES = {
    "Chatbot HS": vista_chatbot_hs,
    "Validador de claims": vista_validador_claims,
    "Generador de material": vista_generador_material,
}

with st.sidebar:
    st.title("Inphormed")
    pagina = st.radio("Menú", list(PAGES.keys()))

PAGES[pagina]()
