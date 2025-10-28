import os
import uuid
from pptx import Presentation
from docx import Document

from app.services.citation_service import extract_citations, normalize
from app.services.fetch_service import fetch_text_for_ref
from app.services.similarity_service import similarity
from app.services.report_service import build_report

# 👇 NUEVOS imports para anotación y snippets
from app.services.annotate_service import (
    annotate_pptx,
    annotate_docx,
    write_snippets_html,
)


class ValidationAgent:
    """
    Orquesta la validación de claims a partir de .pptx o .docx.
    - Extrae texto (slides o párrafos)
    - Detecta referencias (DOI / PMID / autor-año)
    - Recupera texto de la fuente (abstract/título)
    - Calcula similitud y color (verde/amarillo/rojo)
    - Genera PDF de informe
    - (NUEVO) Genera documento anotado + snippets HTML y añade links
    """

    def validate_file(self, file_path: str, out_pdf: str):
        # Asegura directorios de trabajo
        os.makedirs("outputs", exist_ok=True)
        os.makedirs("uploads", exist_ok=True)

        ext = os.path.splitext(file_path)[1].lower()
        findings = []

        # ---------- EXTRACCIÓN DE TEXTO Y POBLADO DE FINDINGS ----------
        if ext == ".pptx":
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    # obtener texto de shapes con propiedad 'text'
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text = "\n".join(t for t in slide_text if t).strip()
                if not text:
                    continue
                findings.extend(self._process_text(text, i))

        elif ext == ".docx":
            doc = Document(file_path)
            for i, para in enumerate(doc.paragraphs, start=1):
                text = para.text.strip()
                if not text:
                    continue
                findings.extend(self._process_text(text, i))

        else:
            raise ValueError("Formato no soportado: usa .pptx o .docx")

        # ---------- ⬇️ AQUI VAN LOS AÑADIDOS ESPECÍFICOS ⬇️ ----------

        # 1) Añadir URLs de origen (paper) por finding (DOI → doi.org, PMID → PubMed)
        #    Si usas resolución previa de autor-año a DOI, puedes ponerla en f['resolved_url'] antes
        for f in findings:
            ref_id = f["ref_raw"]
            # DOI
            if isinstance(ref_id, str) and ref_id.startswith("10."):
                f["source_url"] = f"https://doi.org/{ref_id}"
            # PMID numérico
            elif isinstance(ref_id, str) and ref_id.isdigit():
                f["source_url"] = f"https://pubmed.ncbi.nlm.nih.gov/{ref_id}/"
            else:
                # autor-año (best effort): si antes resolviste algo, respétalo
                f["source_url"] = f.get("resolved_url", "")

        # 2) Generar HTML de snippets y asignar un anchor por claim
        base_uid = uuid.uuid4().hex
        snippets_path = write_snippets_html(base_uid, findings)
        for idx, f in enumerate(findings, start=1):
            # Nota: montas /outputs en main.py, así que este path será accesible como "/outputs/..."
            f["snippet_url"] = f"/{snippets_path}#claim-{idx}"

        # 3) Generar documento anotado (colores + enlaces a snippet/paper)
        annotated_path = None
        if ext == ".pptx":
            annotated_path = f"outputs/annotated_{base_uid}.pptx"
            annotate_pptx(file_path, annotated_path, findings)
        elif ext == ".docx":
            annotated_path = f"outputs/annotated_{base_uid}.docx"
            annotate_docx(file_path, annotated_path, findings)

        # ---------- ⬆️ FIN DE LOS AÑADIDOS ESPECÍFICOS ⬆️ ----------

        # 4) Generar el PDF de informe de colores (como ya hacías)
        build_report(file_path, findings, out_pdf)

        # 5) Resumen agregado para devolver al frontend
        summary = {
            "green": sum(1 for f in findings if f["color"] == "green"),
            "yellow": sum(1 for f in findings if f["color"] == "yellow"),
            "red": sum(1 for f in findings if f["color"] == "red"),
            "total": len(findings),
            "report": out_pdf,
            # rutas añadidas
            "annotated_doc": annotated_path,         # /outputs/annotated_*.pptx|docx
            "snippets_html": snippets_path,          # /outputs/snippets/snippets_*.html
        }
        return summary

    # -----------------------------------------------------------------

    def _process_text(self, text: str, index: int):
        """
        Detecta citas en 'text', recupera fuente, calcula similitud y devuelve hallazgos.
        """
        results = []
        refs = [normalize(r) for r in extract_citations(text)]
        for ref in refs:
            src = fetch_text_for_ref(ref)  # abstract o título+abstract
            score, color = similarity(text, src or "")
            results.append({
                "slide": index,                       # index = nº de slide o párrafo
                "claim": text,                        # texto del bloque
                "ref_raw": ref["id"],                 # DOI/PMID/autor-año
                "score": score,
                "color": color,                       # "green" | "yellow" | "red"
                "source_excerpt": (src or "")[:1200], # recorte de la fuente
                # "source_url" y "snippet_url" se rellenan después
            })
        return results
